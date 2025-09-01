# -*- coding: utf-8 -*-
"""
文档处理服务
支持多种文档格式的解析和文本提取
"""
import os
import io
import logging
from typing import Dict, Any, Optional, Tuple
from werkzeug.datastructures import FileStorage

# 文档解析库
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import docx2txt
    DOC_AVAILABLE = True
except ImportError:
    DOC_AVAILABLE = False

try:
    import subprocess
    import tempfile
    ANTIWORD_AVAILABLE = True
except ImportError:
    ANTIWORD_AVAILABLE = False

try:
    import openpyxl
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import xlrd
    XLS_AVAILABLE = True
except ImportError:
    XLS_AVAILABLE = False

import re
from collections import Counter

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentContentCleaner:
    """文档内容清洗器 - 专门处理OCR噪声和格式化标记"""
    
    def __init__(self):
        # Word文档特有的格式化标记模式
        self.word_noise_patterns = [
            r'HYPERLINK\s+"[^"]*"',  # HYPERLINK "url" 格式
            r'HYPERLINK\s+[^\s]+',   # HYPERLINK url 格式
            r'\bHYPERLINK\b.*?(?=\s|$)',  # 任何包含HYPERLINK的行
            r'\bEMBED\b.*?(?=\s|$)',      # EMBED 标记
            r'\bMERGEFORMAT\b',           # MERGEFORMAT 标记  
            r'\bCERTIFICATE\b',           # CERTIFICATE 标记
            r'\bPACKING\b',               # PACKING 标记
            r'\b_GoBack\b',               # Word书签标记
            r'\b_Toc\d+\b',               # 目录标记
        ]
        
        # 页面和导航元素模式
        self.page_nav_patterns = [
            r'PAGE\s+\d+',                # PAGE 7 格式
            r'第\s*\d+\s*页',             # 第7页 格式
            r'页\s*码\s*[:：]\s*\d+',      # 页码：7 格式
            r'CHAPTER\s+\d+',             # CHAPTER 标记
            r'CONTENTS?(?:\s|$)',         # CONTENTS 目录标记
            r'INDEX(?:\s|$)',             # INDEX 索引标记
            r'TITLE(?:\s|$)',             # TITLE 标记
        ]
        
        # 表格边框和格式化符号
        self.table_noise_patterns = [
            r'^[\|\s\-\+\=]{3,}$',        # 表格边框线
            r'^[A-Za-z]\s+[A-Za-z]\s+[A-Za-z].*[A-Za-z]\s+[A-Za-z]$',  # "A A AB X B" 类型
            r'^[A-Za-z]{1,2}\s+[a-z]\s+[a-z].*[a-z]\s+[a-z]$',        # "Ca a a a b" 类型
            r'^[a-z]+\s+\d+\s+[A-Z]+.*$',  # "h 9 HYPERLINK" 类型
            r'^\s*[｜\|]\s*$',            # 单独的管道符
            r'^\s*[\-]{2,}\s*$',          # 连字符分隔线
        ]
        
        # 无意义字符重复模式
        self.meaningless_patterns = [
            r'^(.)\1{5,}$',               # 同一字符重复5次以上
            r'^\s*[\.]{3,}\s*$',          # 点号省略符
            r'^\s*[\*]{3,}\s*$',          # 星号装饰符
            r'^\s*[\-]{3,}\s*$',          # 连字符分隔
            r'^\s*[_]{3,}\s*$',           # 下划线分隔
        ]
        
        # 技术规格保护模式 - 这些模式的内容需要保留
        self.protect_patterns = [
            r'\d+[VvAaWwHh℃℉%]',         # 包含技术单位
            r'\d+\s*[-~±]\s*\d+',        # 数值范围
            r'\d+\s*[:/]\s*\d+',         # 比值格式
            r'(?:电|压|流|功|率|频|温|度|精|量)', # 电力技术关键词
            r'\b(?:volt|amp|watt|freq|temp|test|spec)\b', # 英文技术词（单词边界）
        ]
    
    def clean_document_content(self, content: str) -> Dict[str, Any]:
        """
        全面清洗文档内容
        
        Args:
            content: 原始文档内容
            
        Returns:
            Dict: 包含清洗结果和统计信息
        """
        original_lines = content.split('\n')
        cleaned_lines = []
        noise_stats = {
            'word_noise': 0,
            'page_nav': 0, 
            'table_noise': 0,
            'meaningless': 0,
            'total_removed': 0,
            'total_lines': len(original_lines)
        }
        
        for line in original_lines:
            line = line.strip()
            if not line:
                continue
                
            # 检查是否为需要保护的技术内容
            if self._is_protected_content(line):
                cleaned_lines.append(line)
                continue
            
            # 检查各种噪声模式
            is_noise = False
            
            # Word格式化噪声
            if self._matches_patterns(line, self.word_noise_patterns):
                noise_stats['word_noise'] += 1
                is_noise = True
            
            # 页面导航元素
            elif self._matches_patterns(line, self.page_nav_patterns):
                noise_stats['page_nav'] += 1
                is_noise = True
            
            # 表格格式化噪声
            elif self._matches_patterns(line, self.table_noise_patterns):
                noise_stats['table_noise'] += 1
                is_noise = True
            
            # 无意义字符重复
            elif self._matches_patterns(line, self.meaningless_patterns):
                noise_stats['meaningless'] += 1
                is_noise = True
            
            if not is_noise:
                # 进一步清理有效行中的噪声标记
                cleaned_line = self._clean_line_noise(line)
                if cleaned_line and len(cleaned_line.strip()) > 1:
                    cleaned_lines.append(cleaned_line)
        
        noise_stats['total_removed'] = len(original_lines) - len(cleaned_lines)
        cleaned_content = '\n'.join(cleaned_lines)
        
        logger.info(f"文档清洗完成: 原始{noise_stats['total_lines']}行 -> 清洗后{len(cleaned_lines)}行, "
                   f"移除噪声{noise_stats['total_removed']}行")
        
        return {
            'cleaned_content': cleaned_content,
            'original_content': content,
            'noise_statistics': noise_stats,
            'cleaning_ratio': noise_stats['total_removed'] / noise_stats['total_lines'] if noise_stats['total_lines'] > 0 else 0
        }
    
    def _is_protected_content(self, line: str) -> bool:
        """检查是否为需要保护的技术内容"""
        return any(re.search(pattern, line, re.IGNORECASE) for pattern in self.protect_patterns)
    
    def _matches_patterns(self, text: str, patterns: list) -> bool:
        """检查文本是否匹配任何噪声模式"""
        return any(re.search(pattern, text, re.IGNORECASE) for pattern in patterns)
    
    def _clean_line_noise(self, line: str) -> str:
        """清理行内的噪声标记"""
        cleaned = line
        
        # 移除Word格式化标记
        for pattern in self.word_noise_patterns:
            cleaned = re.sub(pattern, '', cleaned, flags=re.IGNORECASE)
        
        # 清理多余空格
        cleaned = re.sub(r'\s+', ' ', cleaned).strip()
        
        return cleaned

class DocumentProcessor:
    """文档处理器"""
    
    # 支持的文档类型
    SUPPORTED_TYPES = {
        # 文本文件
        'text/plain': 'txt',
        
        # PDF文件
        'application/pdf': 'pdf',
        
        # Word文档
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/msword': 'doc',
        
        # Excel文件
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'application/vnd.ms-excel': 'xls',
        
        # PowerPoint文件
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-powerpoint': 'ppt',
        
        # RTF文件
        'application/rtf': 'rtf',
        'text/rtf': 'rtf',
        
        # 图片文件
        'image/png': 'png',
        'image/jpeg': 'jpg',
        'image/jpg': 'jpg',
        'image/gif': 'gif',
        'image/bmp': 'bmp',
        'image/tiff': 'tiff',
    }
    
    def __init__(self):
        self.max_file_size = 10 * 1024 * 1024  # 10MB
        self.max_text_length = 100000  # 增加到100k字符限制，支持更长的技术文档
        self.min_text_length = 10  # 最小文本长度要求
        self.content_cleaner = DocumentContentCleaner()  # 文档内容清洗器
    
    def is_supported_format(self, mimetype: str, filename: str = "") -> bool:
        """检查文件格式是否支持"""
        if mimetype in self.SUPPORTED_TYPES:
            return True
        
        # 基于文件扩展名检查
        if filename:
            ext = filename.lower().split('.')[-1]
            return ext in ['txt', 'pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'rtf', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']
        
        return False
    
    def get_supported_formats(self) -> Dict[str, list]:
        """获取支持的格式信息"""
        return {
            'mimetypes': list(self.SUPPORTED_TYPES.keys()),
            'extensions': ['txt', 'pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'rtf', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff'],
            'availability': {
                'pdf': PDF_AVAILABLE,
                'docx': DOCX_AVAILABLE,
                'doc': DOC_AVAILABLE,
                'xlsx': XLSX_AVAILABLE,
                'xls': XLS_AVAILABLE,
                'pptx': PPTX_AVAILABLE,
                'rtf': RTF_AVAILABLE,
                'ocr': OCR_AVAILABLE
            }
        }
    
    def process_document(self, file: FileStorage) -> Tuple[str, Dict[str, Any]]:
        """
        处理上传的文档
        
        Args:
            file: 上传的文件对象
            
        Returns:
            Tuple[str, Dict]: (提取的文本内容, 文档信息)
        """
        # 验证文件
        validation_result = self._validate_file(file)
        if not validation_result['valid']:
            raise ValueError(validation_result['error'])
        
        # 获取文件信息
        doc_info = {
            'filename': file.filename,
            'mimetype': file.mimetype,
            'size': self._get_file_size(file),
            'type': self._detect_file_type(file.mimetype, file.filename)
        }
        
        try:
            # 根据文件类型提取文本
            raw_text_content = self._extract_text(file, doc_info['type'])
            
            # 🔧 文档内容清洗 - 移除OCR噪声和格式化标记
            cleaning_result = self.content_cleaner.clean_document_content(raw_text_content)
            text_content = cleaning_result['cleaned_content']
            
            # 记录清洗统计信息
            doc_info['cleaning_stats'] = cleaning_result['noise_statistics']
            doc_info['cleaning_ratio'] = cleaning_result['cleaning_ratio']
            
            # 智能文本长度处理
            if len(text_content) > self.max_text_length:
                # 使用智能截断，保留重要信息
                text_content = self._intelligent_truncate(text_content)
                doc_info['truncated'] = True
                logger.warning(f"Text content intelligently truncated to {len(text_content)} characters")
            
            doc_info['text_length'] = len(text_content)
            doc_info['word_count'] = len(text_content.split())
            doc_info['original_text_length'] = len(raw_text_content)
            
            # 日志记录清洗效果
            if cleaning_result['cleaning_ratio'] > 0.1:  # 清洗比例超过10%时记录
                logger.info(f"文档清洗效果显著: 清洗比例{cleaning_result['cleaning_ratio']:.1%}, "
                          f"移除{cleaning_result['noise_statistics']['total_removed']}行噪声")
            
            return text_content, doc_info
            
        except Exception as e:
            logger.error(f"Error processing document {file.filename}: {str(e)}")
            raise ValueError(f"Document processing failed: {str(e)}")
    
    def _validate_file(self, file: FileStorage) -> Dict[str, Any]:
        """验证上传文件"""
        if not file:
            return {'valid': False, 'error': 'No file provided'}
        
        if not file.filename:
            return {'valid': False, 'error': 'No filename provided'}
        
        # 检查文件大小
        file_size = self._get_file_size(file)
        if file_size > self.max_file_size:
            return {
                'valid': False, 
                'error': f'File size ({file_size / 1024 / 1024:.1f}MB) exceeds limit ({self.max_file_size / 1024 / 1024}MB)'
            }
        
        # 检查文件格式
        if not self.is_supported_format(file.mimetype or '', file.filename):
            return {
                'valid': False,
                'error': f'Unsupported file format: {file.mimetype or "unknown"}'
            }
        
        return {'valid': True}
    
    def _get_file_size(self, file: FileStorage) -> int:
        """获取文件大小"""
        current_pos = file.tell()
        file.seek(0, os.SEEK_END)
        size = file.tell()
        file.seek(current_pos)
        return size
    
    def _detect_file_type(self, mimetype: str, filename: str) -> str:
        """检测文件类型"""
        if mimetype in self.SUPPORTED_TYPES:
            return self.SUPPORTED_TYPES[mimetype]
        
        # 基于扩展名检测
        if filename:
            ext = filename.lower().split('.')[-1]
            if ext in ['txt']:
                return 'txt'
            elif ext in ['pdf']:
                return 'pdf'
            elif ext in ['docx']:
                return 'docx'
            elif ext in ['doc']:
                return 'doc'
            elif ext in ['xlsx']:
                return 'xlsx'
            elif ext in ['xls']:
                return 'xls'
            elif ext in ['pptx']:
                return 'pptx'
            elif ext in ['ppt']:
                return 'ppt'
            elif ext in ['rtf']:
                return 'rtf'
            elif ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
                return ext
        
        return 'unknown'
    
    def _extract_text(self, file: FileStorage, file_type: str) -> str:
        """根据文件类型提取文本内容"""
        file.seek(0)  # 重置文件指针
        
        if file_type == 'txt':
            return self._extract_text_from_txt(file)
        elif file_type == 'pdf':
            return self._extract_text_from_pdf(file)
        elif file_type == 'docx':
            return self._extract_text_from_docx(file)
        elif file_type == 'doc':
            return self._extract_text_from_doc(file)
        elif file_type == 'xlsx':
            return self._extract_text_from_xlsx(file)
        elif file_type == 'xls':
            return self._extract_text_from_xls(file)
        elif file_type in ['pptx', 'ppt']:
            return self._extract_text_from_pptx(file)
        elif file_type == 'rtf':
            return self._extract_text_from_rtf(file)
        elif file_type in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
            return self._extract_text_from_image(file)
        else:
            raise ValueError(f"Unsupported file type for text extraction: {file_type}")
    
    def _extract_text_from_txt(self, file: FileStorage) -> str:
        """从TXT文件提取文本"""
        try:
            content = file.read()
            # 尝试不同编码
            for encoding in ['utf-8', 'gbk', 'gb2312', 'latin1']:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # 如果所有编码都失败，使用错误处理
            return content.decode('utf-8', errors='ignore')
        except Exception as e:
            raise ValueError(f"Failed to read text file: {str(e)}")
    
    def _extract_text_from_pdf(self, file: FileStorage) -> str:
        """从PDF文件提取文本"""
        if not PDF_AVAILABLE:
            raise ValueError("PDF processing not available. Please install PyPDF2.")
        
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
            text_content = ""
            
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
            
            if not text_content.strip():
                raise ValueError("No text content found in PDF. The PDF might be image-based.")
            
            return text_content
            
        except Exception as e:
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")
    
    def _extract_text_from_docx(self, file: FileStorage) -> str:
        """从DOCX文件提取文本"""
        if not DOCX_AVAILABLE:
            raise ValueError("DOCX processing not available. Please install python-docx.")
        
        try:
            doc = Document(io.BytesIO(file.read()))
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + "\t"
                    text_content += "\n"
            
            return text_content
            
        except Exception as e:
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")
    
    def _extract_text_from_image(self, file: FileStorage) -> str:
        """从图片文件提取文本(OCR)"""
        if not OCR_AVAILABLE:
            raise ValueError("OCR processing not available. Please install pytesseract and Pillow.")
        
        try:
            image = Image.open(io.BytesIO(file.read()))
            
            # 🔧 增强OCR处理 - 尝试多种OCR配置
            ocr_configs = [
                '--psm 6',    # 统一文本块
                '--psm 4',    # 假设为单列文本
                '--psm 3',    # 全自动页面分割但不进行方向和脚本检测
                '--psm 8',    # 将图像视为单个单词
                '--psm 13'    # 原始行。将图像视为单个文本行
            ]
            
            best_text = ""
            best_score = 0
            
            for config in ocr_configs:
                try:
                    # 使用tesseract进行OCR，支持中英文
                    text_content = pytesseract.image_to_string(
                        image, 
                        lang='chi_sim+eng',  # 中文简体+英文
                        config=config
                    )
                    
                    if text_content and text_content.strip():
                        # 评估OCR结果质量
                        quality_score = self._evaluate_ocr_quality(text_content)
                        logger.info(f"OCR配置 {config} 质量评分: {quality_score}")
                        
                        if quality_score > best_score:
                            best_text = text_content
                            best_score = quality_score
                            
                        # 如果质量很高，提前退出
                        if quality_score > 0.8:
                            break
                            
                except Exception as ocr_error:
                    logger.debug(f"OCR配置 {config} 失败: {str(ocr_error)}")
                    continue
            
            # 如果没有找到任何文本，尝试英文专用OCR
            if not best_text.strip():
                try:
                    text_content = pytesseract.image_to_string(
                        image,
                        lang='eng',  # 仅英文
                        config='--psm 6'
                    )
                    if text_content and text_content.strip():
                        quality_score = self._evaluate_ocr_quality(text_content)
                        if quality_score > best_score:
                            best_text = text_content
                            best_score = quality_score
                            logger.info(f"英文OCR质量评分: {quality_score}")
                except Exception as e:
                    logger.debug(f"英文OCR失败: {str(e)}")
            
            # 清理OCR结果
            if best_text.strip():
                cleaned_text = self._clean_ocr_text(best_text)
                if len(cleaned_text.strip()) >= 10:  # 降低最小长度要求
                    logger.info(f"OCR成功提取文本，质量评分: {best_score:.2f}")
                    return cleaned_text
            
            # 如果所有方法都失败，提供详细错误信息
            raise ValueError(
                "无法从图片中提取可读文本。可能的原因：\n"
                "1. 图片中没有文本内容\n"
                "2. 图片分辨率过低或质量较差\n"
                "3. 文字颜色与背景对比度不足\n"
                "4. 文字角度倾斜或方向不正确\n\n"
                "建议解决方案：\n"
                "• 提供更高分辨率的图片\n"
                "• 确保文字清晰可读\n"
                "• 将图片转换为文档格式（PDF、Word等）"
            )
            
        except Exception as e:
            if "无法从图片中提取可读文本" in str(e):
                raise e
            else:
                raise ValueError(f"图片处理失败: {str(e)}")
    
    def _extract_text_from_doc(self, file: FileStorage) -> str:
        """从DOC文件提取文本 - 使用多种方法尝试"""
        file_content = file.read()
        best_text = ""
        best_score = 0
        extraction_methods = []
        
        # 方法1: 尝试使用native antiword工具 (如果存在)
        try:
            with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file.flush()
                
                try:
                    # 尝试系统级antiword命令
                    result = subprocess.run(
                        ['/usr/local/bin/antiword', temp_file.name],
                        capture_output=True,
                        text=True,
                        timeout=30
                    )
                    
                    if result.returncode == 0 and result.stdout.strip():
                        quality_score = self._evaluate_text_quality(result.stdout)
                        extraction_methods.append(("native-antiword", result.stdout, quality_score))
                        logger.info(f"Native antiword extraction quality score: {quality_score}")
                        if quality_score > best_score:
                            best_text = result.stdout
                            best_score = quality_score
                    else:
                        logger.debug(f"Native antiword failed with return code {result.returncode}")
                        
                except (subprocess.TimeoutExpired, FileNotFoundError) as e:
                    logger.debug(f"Native antiword command not found: {str(e)}")
                finally:
                    # 清理临时文件
                    try:
                        os.unlink(temp_file.name)
                    except:
                        pass
                        
        except Exception as e:
            logger.debug(f"Native antiword method failed: {str(e)}")
        
        # 方法2: 尝试使用docx2txt (适用于一些.doc文件)
        if DOC_AVAILABLE:
            try:
                file.seek(0)
                text_content = docx2txt.process(io.BytesIO(file_content))
                if text_content and text_content.strip():
                    quality_score = self._evaluate_text_quality(text_content)
                    extraction_methods.append(("docx2txt", text_content, quality_score))
                    logger.info(f"docx2txt extraction quality score: {quality_score}")
                    if quality_score > best_score:
                        best_text = text_content
                        best_score = quality_score
            except Exception as e:
                logger.debug(f"docx2txt failed: {str(e)}")
        
        # 方法3: OLE文件结构解析 (新增方法)
        try:
            import olefile
            
            # 检查是否是有效的OLE文件
            if olefile.isOleFile(io.BytesIO(file_content)):
                with olefile.OleFileIO(io.BytesIO(file_content)) as ole:
                    # 尝试读取WordDocument流
                    try:
                        # 检查WordDocument流是否存在
                        ole.openstream('WordDocument')
                        ole_has_word_stream = True
                    except:
                        ole_has_word_stream = False
                    
                    if ole_has_word_stream:
                        try:
                            word_stream = ole.openstream('WordDocument')
                            raw_data = word_stream.read()
                            
                            # 尝试多种编码解析
                            encodings_to_try = ['utf-16le', 'utf-16', 'gbk', 'gb2312', 'utf-8', 'latin1']
                            
                            for encoding in encodings_to_try:
                                try:
                                    decoded_text = raw_data.decode(encoding, errors='ignore')
                                    clean_text = self._clean_ole_text(decoded_text)
                                    
                                    if len(clean_text.strip()) > 100:  # 至少需要100个字符
                                        quality_score = self._evaluate_text_quality(clean_text)
                                        extraction_methods.append((f"ole-{encoding}", clean_text, quality_score))
                                        logger.info(f"OLE {encoding} extraction quality score: {quality_score}")
                                        if quality_score > best_score:
                                            best_text = clean_text
                                            best_score = quality_score
                                except:
                                    continue
                                    
                        except Exception as e:
                            logger.debug(f"WordDocument stream extraction failed: {e}")
                            
        except Exception as e:
            logger.debug(f"OLE extraction method failed: {str(e)}")
        
        # 方法4: 增强编码检测和转换
        # 🔧 使用更多编码方式和检测策略
        encodings_to_try = [
            # 中文编码 - 优先尝试
            'gb18030', 'gbk', 'gb2312', 'big5', 'hz',
            # Unicode编码
            'utf-8', 'utf-16', 'utf-16le', 'utf-16be', 'utf-32',
            # Windows编码
            'cp936', 'cp950', 'cp1252', 'cp1251',
            # ISO编码
            'iso-8859-1', 'latin1'
        ]
        
        # 🔧 增强编码检测 - 先尝试检测最可能的编码
        try:
            import chardet
            
            # 检测文件编码
            file.seek(0)
            sample = file.read(min(10240, len(file_content)))  # 读取前10KB检测编码
            encoding_result = chardet.detect(sample)
            
            if encoding_result and encoding_result.get('encoding') and encoding_result.get('confidence', 0) > 0.7:
                detected_encoding = encoding_result['encoding'].lower()
                logger.info(f"检测到编码: {detected_encoding} (置信度: {encoding_result['confidence']})")
                
                # 将检测到的编码放在列表首位
                if detected_encoding not in [enc.lower() for enc in encodings_to_try]:
                    encodings_to_try.insert(0, detected_encoding)
                else:
                    # 将检测到的编码移到前面
                    for i, enc in enumerate(encodings_to_try):
                        if enc.lower() == detected_encoding:
                            encodings_to_try.insert(0, encodings_to_try.pop(i))
                            break
                            
        except ImportError:
            logger.debug("chardet模块未安装，跳过自动编码检测")
        except Exception as e:
            logger.debug(f"编码检测失败: {str(e)}")
        
        # 尝试不同编码
        for encoding in encodings_to_try:
            try:
                text_content = file_content.decode(encoding, errors='replace')  # 使用replace而不是ignore
                
                # 🔧 更严格的文本验证
                if self._is_corrupted_text(text_content):
                    logger.debug(f"编码 {encoding} 产生乱码，跳过")
                    continue
                
                # 高级文本清理
                clean_text = self._clean_extracted_text(text_content)
                if len(clean_text.strip()) > self.min_text_length:
                    quality_score = self._evaluate_text_quality(clean_text)
                    extraction_methods.append((f"encoding-{encoding}", clean_text, quality_score))
                    logger.debug(f"Encoding {encoding} extraction quality score: {quality_score}")
                    if quality_score > best_score:
                        best_text = clean_text
                        best_score = quality_score
                        
                    # 如果质量分数很高，提前退出
                    if quality_score > 0.8:
                        logger.info(f"找到高质量文本 (质量分数: {quality_score})，使用编码: {encoding}")
                        break
                        
            except Exception as e:
                logger.debug(f"编码 {encoding} 解码失败: {str(e)}")
                continue
        
        # 记录所有尝试的方法
        logger.info(f"DOC extraction attempted {len(extraction_methods)} methods, best score: {best_score}")
        
        # 如果找到了可用的文本
        if best_text and best_score > 0.1:  # 最低质量阈值
            logger.info(f"Successfully extracted DOC text with quality score {best_score}")
            return best_text
        
        # 如果所有方法都失败，提供更详细的错误信息和建议
        error_msg = (
            "无法从此.doc文件中提取可读文本。可能的原因："
            "\n1. 文件损坏或格式不标准"
            "\n2. 文件包含主要是图片或图表"
            "\n3. 文件受密码保护"
            "\n4. 文件使用了特殊的编码格式"
            "\n\n建议解决方案："
            "\n• 请尝试将文件转换为.docx格式"
            "\n• 确保文件包含可读的文本内容"
            "\n• 检查文件是否完整下载"
            f"\n• 尝试的提取方法数量: {len(extraction_methods)}"
        )
        
        raise ValueError(error_msg)
    
    def _extract_text_from_xlsx(self, file: FileStorage) -> str:
        """从XLSX文件提取文本"""
        if not XLSX_AVAILABLE:
            raise ValueError("XLSX processing not available. Please install openpyxl.")
        
        try:
            workbook = load_workbook(io.BytesIO(file.read()), read_only=True, data_only=True)
            text_content = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"=== Sheet: {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text_content += "\t".join(row_text) + "\n"
                
                text_content += "\n"
            
            workbook.close()
            
            if not text_content.strip():
                raise ValueError("No text content found in XLSX file.")
            
            return text_content
            
        except Exception as e:
            raise ValueError(f"Failed to extract text from XLSX: {str(e)}")
    
    def _extract_text_from_xls(self, file: FileStorage) -> str:
        """从XLS文件提取文本"""
        if not XLS_AVAILABLE:
            raise ValueError("XLS processing not available. Please install xlrd.")
        
        try:
            workbook = xlrd.open_workbook(file_contents=file.read())
            text_content = ""
            
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text_content += f"=== Sheet: {sheet.name} ===\n"
                
                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            row_values.append(str(cell.value))
                    if row_values:
                        text_content += "\t".join(row_values) + "\n"
                
                text_content += "\n"
            
            if not text_content.strip():
                raise ValueError("No text content found in XLS file.")
            
            return text_content
            
        except Exception as e:
            raise ValueError(f"Failed to extract text from XLS: {str(e)}")
    
    def _extract_text_from_pptx(self, file: FileStorage) -> str:
        """从PPTX/PPT文件提取文本"""
        if not PPTX_AVAILABLE:
            raise ValueError("PPTX processing not available. Please install python-pptx.")
        
        try:
            presentation = Presentation(io.BytesIO(file.read()))
            text_content = ""
            
            for slide_idx, slide in enumerate(presentation.slides):
                text_content += f"=== Slide {slide_idx + 1} ===\n"
                
                # 提取幻灯片中的所有文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
                    elif hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_content += run.text
                            text_content += "\n"
                
                text_content += "\n"
            
            if not text_content.strip():
                raise ValueError("No text content found in PPTX file.")
            
            return text_content
            
        except Exception as e:
            raise ValueError(f"Failed to extract text from PPTX: {str(e)}")
    
    def _extract_text_from_rtf(self, file: FileStorage) -> str:
        """从RTF文件提取文本"""
        if not RTF_AVAILABLE:
            raise ValueError("RTF processing not available. Please install striprtf.")
        
        try:
            rtf_content = file.read().decode('utf-8', errors='ignore')
            text_content = rtf_to_text(rtf_content)
            
            if not text_content or not text_content.strip():
                raise ValueError("No text content found in RTF file.")
            
            return text_content
            
        except Exception as e:
            # 尝试其他编码
            try:
                file.seek(0)
                rtf_content = file.read().decode('latin1', errors='ignore')
                text_content = rtf_to_text(rtf_content)
                
                if text_content and text_content.strip():
                    return text_content
                else:
                    raise ValueError("No text content found in RTF file.")
                    
            except Exception as e2:
                raise ValueError(f"Failed to extract text from RTF: {str(e)} / {str(e2)}")
    
    def _clean_extracted_text(self, text: str) -> str:
        """清理从文档中提取的文本 - 增强版，专门处理.doc文件的二进制垃圾数据"""
        import re
        
        # 1. 🛡️ 高级二进制垃圾过滤 - 移除明显的编码错误字符
        # 移除控制字符和非打印字符（保留换行符、制表符和空格）
        cleaned = ''.join(char for char in text if char.isprintable() or char in '\n\t ')
        
        # 🔧 移除.doc解析常见的乱码字符模式
        doc_garbage_patterns = [
            r'[潗摲楍牣獯景煅慵楴湯畱瑡潩卍潗摲潄吀瑩敬牁慩袈霡蠈袢]+',  # 常见.doc乱码字符串
            r'[㸳㠴㔷㤸㜹㈰㐱㠲㌳㘴㔵㘶㠷㤸㠹]+',  # 十六进制乱码
            r'[屜屝屬屭屨屪屢屣層履屦屧屨屩屲]+',  # Word结构字符
            r'[▉▊▋▌▍▎▏█▄▀■□▲△▼▽◆◇○●◎☆★]+',  # OLE图形字符
            r'[\u0080-\u00ff]{2,}',              # Latin-1扩展字符乱码
            r'[\ue000-\uf8ff]+',                 # 私用区字符
            r'[\ufeff\ufffe]+',                  # 字节序标记
        ]
        
        for pattern in doc_garbage_patterns:
            cleaned = re.sub(pattern, ' ', cleaned)
            
        # 🔧 移除编码错误产生的替代字符和问号串
        cleaned = re.sub(r'[�?]{2,}', ' ', cleaned)
        
        # 2. 🔍 智能文本分割和识别
        # 使用改进的正则表达式查找中文文本块
        chinese_pattern = r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+'
        chinese_matches = re.findall(chinese_pattern, cleaned)
        
        # 查找英文文本块（更严格的模式，排除乱码）
        english_pattern = r'[a-zA-Z][a-zA-Z0-9\s\.,;:!?\-()\'"\[\]{}]{3,}'
        english_matches = re.findall(english_pattern, cleaned)
        
        # 查找技术数据块（数字+单位的模式）
        tech_data_pattern = r'\d+[\w\s]*[A-Za-z%°℃Ω/\-]{1,5}'
        tech_matches = re.findall(tech_data_pattern, cleaned)
        
        # 3. 🧹 智能文本重组
        meaningful_text = []
        
        # 添加中文文本（更严格的过滤）
        for match in chinese_matches:
            clean_match = match.strip()
            # 检查是否包含足够的汉字内容
            chinese_char_count = len(re.findall(r'[\u4e00-\u9fff]', clean_match))
            if chinese_char_count >= 2 and len(clean_match) >= 3:
                # 额外检查：确保不是乱码汉字
                if not self._contains_garbage_chinese(clean_match):
                    meaningful_text.append(clean_match)
        
        # 添加英文文本（过滤掉可能的乱码）
        for match in english_matches:
            clean_match = match.strip()
            words = clean_match.split()
            # 更严格的英文内容验证
            if (len(words) >= 2 and len(clean_match) >= 8 and 
                self._is_meaningful_english(clean_match)):
                meaningful_text.append(clean_match)
        
        # 添加技术数据
        for match in tech_matches:
            clean_match = match.strip()
            if len(clean_match) >= 3 and not self._is_corrupted_tech_data(clean_match):
                meaningful_text.append(clean_match)
        
        # 4. 🔄 如果智能提取失败，使用保守的基本清理
        if not meaningful_text:
            logger.warning("智能文本提取未找到有意义内容，使用保守清理")
            # 基本清理：移除连续的空白字符
            cleaned = re.sub(r'\s+', ' ', cleaned)
            # 移除明显的垃圾行
            lines = []
            for line in cleaned.split('\n'):
                line = line.strip()
                if (line and len(line) > 2 and 
                    not self._is_obvious_garbage_line(line)):
                    lines.append(line)
            return '\n'.join(lines)
        
        # 5. 🎯 组合有意义的文本
        result = '\n'.join(meaningful_text)
        
        # 6. 🧹 最终清理和规范化
        result = re.sub(r'\n\s*\n', '\n\n', result)  # 规范段落间距
        result = re.sub(r'[ \t]+', ' ', result)      # 规范空格
        result = re.sub(r'\n{3,}', '\n\n', result)   # 限制最多双换行
        
        return result.strip()
    
    def _clean_ole_text(self, text: str) -> str:
        """专门清理从OLE文件结构中提取的文本"""
        import re
        
        # 1. 移除所有控制字符和非打印字符
        cleaned = ''.join(char for char in text if char.isprintable() or char in '\n\t ')
        
        # 2. 移除Word文档特有的结构化数据
        # 移除OLE对象标识符
        cleaned = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', cleaned)
        
        # 3. 查找并提取有意义的文本段落
        # 寻找中文字符集群
        chinese_pattern = r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]{2,}'
        chinese_blocks = re.findall(chinese_pattern, cleaned)
        
        # 寻找英文字符集群
        english_pattern = r'[a-zA-Z][a-zA-Z0-9\s\.,;:!?\-()\'"\[\]{}]{5,}'
        english_blocks = re.findall(english_pattern, cleaned)
        
        # 4. 过滤和重新组织文本
        meaningful_blocks = []
        
        # 处理中文文本块
        for block in chinese_blocks:
            clean_block = block.strip()
            if len(clean_block) >= 3:  # 至少3个字符
                meaningful_blocks.append(clean_block)
        
        # 处理英文文本块
        for block in english_blocks:
            clean_block = block.strip()
            words = clean_block.split()
            if len(words) >= 3 and len(clean_block) >= 10:  # 至少3个单词且10个字符
                meaningful_blocks.append(clean_block)
        
        # 5. 如果没有找到有意义的内容，回退到基础清理
        if not meaningful_blocks:
            # 移除重复的空白字符
            cleaned = re.sub(r'\s+', ' ', cleaned)
            # 移除明显的乱码模式
            cleaned = re.sub(r'[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffef\w\s\.,;:!?\-()\'"\[\]{}]', '', cleaned)
            return cleaned.strip()
        
        # 6. 组合有意义的文本块
        result = '\n'.join(meaningful_blocks)
        
        # 7. 最终清理
        result = re.sub(r'\n+', '\n\n', result)  # 规范段落间距
        result = re.sub(r'[ \t]+', ' ', result)  # 规范空格
        
        return result.strip()
    
    def _intelligent_truncate(self, text: str) -> str:
        """智能截断文本，保留重要信息"""
        if len(text) <= self.max_text_length:
            return text
        
        # 分割成段落
        paragraphs = text.split('\n\n')
        
        # 优先保留包含关键信息的段落
        important_keywords = [
            # 中文关键词
            '产品', '型号', '规格', '参数', '技术', '电压', '电流', '功率', '频率', 
            '测试', '保护', '继电', '装置', '设备', '性能', '精度', '范围', 
            '温度', '湿度', '环境', '认证', '标准', '质量', '保修', '服务',
            '配置', '接口', '通信', '控制', '显示', '操作', '安装', '维护',
            
            # 英文关键词
            'product', 'model', 'specification', 'parameter', 'technical', 
            'voltage', 'current', 'power', 'frequency', 'test', 'protection',
            'relay', 'device', 'equipment', 'performance', 'accuracy', 'range',
            'temperature', 'humidity', 'environment', 'certification', 'standard',
            'quality', 'warranty', 'service', 'configuration', 'interface',
            'communication', 'control', 'display', 'operation', 'installation'
        ]
        
        # 按重要性对段落排序
        scored_paragraphs = []
        for para in paragraphs:
            score = 0
            para_lower = para.lower()
            
            # 计算关键词得分
            for keyword in important_keywords:
                score += para_lower.count(keyword.lower()) * 10
            
            # 包含数字和单位的段落更重要（技术规格）
            import re
            if re.search(r'\d+[A-Za-z]*[\s]*[A-Za-z/Ω%°℃]', para):
                score += 20
            
            # 包含表格结构的段落
            if '\t' in para or para.count('|') > 2:
                score += 15
            
            # 段落长度适中的更重要
            if 50 <= len(para) <= 500:
                score += 5
            
            scored_paragraphs.append((score, para))
        
        # 按得分排序，保留高分段落
        scored_paragraphs.sort(key=lambda x: x[0], reverse=True)
        
        result = []
        current_length = 0
        target_length = int(self.max_text_length * 0.95)  # 留一些余量
        
        # 首先添加最重要的段落
        for score, para in scored_paragraphs:
            if current_length + len(para) + 2 <= target_length:  # +2 for \n\n
                result.append(para)
                current_length += len(para) + 2
            elif current_length < target_length * 0.7:  # 如果还没达到70%，尝试部分截取
                remaining = target_length - current_length - 2
                if remaining > 100:  # 只有足够长度才截取
                    truncated_para = para[:remaining] + "..."
                    result.append(truncated_para)
                    break
            else:
                break
        
        final_text = '\n\n'.join(result)
        
        # 如果结果太短，从原文开头补充
        if len(final_text) < target_length * 0.5:
            remaining_space = target_length - len(final_text)
            if remaining_space > 100:
                prefix = text[:remaining_space-len(final_text)-20] + "...\n\n"
                final_text = prefix + final_text
        
        logger.info(f"智能截断：原文{len(text)}字符 -> {len(final_text)}字符，保留{len(result)}个段落")
        return final_text
    
    def _is_corrupted_text(self, text: str) -> bool:
        """检测文本是否为乱码"""
        if not text or len(text.strip()) < 20:
            return True
        
        # 🔍 检测乱码的几个指标
        text_sample = text[:2000]  # 取前2000字符作为样本
        
        # 1. 检测异常高频的特殊字符
        special_chars = sum(1 for c in text_sample if ord(c) > 0xFFFF or c in '揀釨娐醢藠俹牁慵楴湯畱瑡潩')
        if special_chars / len(text_sample) > 0.3:  # 超过30%的异常字符
            logger.debug(f"检测到高频异常字符: {special_chars}/{len(text_sample)}")
            return True
        
        # 2. 检测重复模式 (常见乱码特征) - 调整阈值，避免过度检测
        import re
        # 检测重复的短字符串模式（如"橢橢"、"塅塅"等）
        repeated_patterns = re.findall(r'(.{1,3})\1{5,}', text_sample)
        if len(repeated_patterns) > 15:  # 提高阈值
            logger.debug(f"检测到大量重复模式: {len(repeated_patterns)}")
            return True
        
        # 3. 检测是否包含可读的中文或英文内容
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text_sample))
        english_words = len(re.findall(r'\b[a-zA-Z]{3,}\b', text_sample))
        readable_content_ratio = (chinese_chars + english_words * 5) / len(text_sample)
        
        if readable_content_ratio < 0.05:  # 可读内容少于5%
            logger.debug(f"检测到可读内容过少: 中文字符{chinese_chars}, 英文单词{english_words}")
            return True
        
        # 4. 检测编码问题特征 - 检查是否有明显的编码错误模式
        encoding_error_patterns = [
            r'[潗摲|楍牣獯景|煅慵楴湯|畱瑡潩|卍潗摲潄|吀瑩敬|牁慩]{5,}',  # 连续的编码错误字符
            r'[袈袈袈|霡蠈袢]{3,}',  # 重复的乱码字符
            r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x84\x86-\x9F]{2,}'  # 控制字符
        ]
        
        for pattern in encoding_error_patterns:
            if re.search(pattern, text_sample):
                logger.debug(f"检测到编码错误模式: {pattern}")
                return True
        
        return False
    
    def _calculate_encoding_score(self, text: str) -> float:
        """计算文本编码质量评分（0-1）"""
        if not text:
            return 0.0
        
        score = 0.0
        
        # 1. 检查UTF-8编码有效性
        try:
            text.encode('utf-8').decode('utf-8')
            score += 0.3
        except UnicodeError:
            pass
        
        # 2. 检查常见编码的兼容性
        for encoding in ['gbk', 'gb2312', 'utf-16']:
            try:
                text.encode(encoding).decode(encoding)
                score += 0.2
                break
            except UnicodeError:
                pass
        
        # 3. 检查控制字符比例
        control_chars = sum(1 for c in text if ord(c) < 32 and c not in '\n\t\r')
        control_ratio = control_chars / len(text)
        if control_ratio < 0.05:  # 控制字符少于5%
            score += 0.3
        
        # 4. 检查字符分布合理性
        printable_chars = sum(1 for c in text if c.isprintable() or c in '\n\t\r')
        printable_ratio = printable_chars / len(text)
        score += printable_ratio * 0.2
        
        return min(score, 1.0)
    
    def _calculate_readability_score(self, text: str) -> float:
        """计算文本可读性评分（0-1）"""
        if not text:
            return 0.0
        
        import re
        
        # 计算各种字符类型
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
        english_words = len(re.findall(r'\b[a-zA-Z]{2,}\b', text))
        digits = len(re.findall(r'\d', text))
        symbols = len(re.findall(r'[\u3000-\u303f\uff00-\uffef]', text))  # 中文标点
        
        # 计算可读内容比例
        readable_content = chinese_chars + english_words * 3 + digits * 0.5 + symbols * 0.8
        readability_ratio = readable_content / len(text)
        
        # 技术内容加成
        technical_bonus = 0
        if re.search(r'(电压|电流|功率|频率|测试|保护|继电|装置|设备|性能|精度|规格|参数|型号)', text):
            technical_bonus += 0.1
        if re.search(r'\d+[A-Za-z%°℃Ω]+', text):  # 数字+单位
            technical_bonus += 0.1
        
        return min(readability_ratio + technical_bonus, 1.0)
    
    def _has_technical_content(self, text: str) -> bool:
        """检查文本是否包含技术内容"""
        import re
        
        # 技术关键词模式
        tech_patterns = [
            r'(电压|电流|功率|频率|电阻|电容|电感)',
            r'(测试|测量|校验|检测|监控)',
            r'(保护|继电|安全|可靠)',
            r'(装置|设备|仪器|仪表)',
            r'(性能|精度|范围|规格)',
            r'(额定|最大|最小|标准)',
            r'\d+[A-Za-z%°℃ΩVAKW]+',  # 数字+单位
            r'[A-Z][A-Z0-9-]{2,}',  # 产品型号
        ]
        
        for pattern in tech_patterns:
            if re.search(pattern, text):
                return True
        
        return False
    
    def _enhanced_corruption_detection(self, text: str, file_type: str = '') -> bool:
        """增强的乱码检测算法"""
        if not text or len(text.strip()) < 20:
            return True
        
        text_sample = text[:3000]  # 扩大检测样本
        
        # 1. 智能编码检测和验证
        encoding_score = self._calculate_encoding_score(text_sample)
        if encoding_score < 0.3:  # 编码质量过低
            logger.debug(f"编码质量评分过低: {encoding_score}")
            return True
        
        # 2. 检测异常高频的特殊字符（精确匹配）
        corruption_chars = sum(1 for c in text_sample if ord(c) > 0xFFFF or c in '揀釨娐醢藠俹牁慵楴湯畱瑡潩卍潗摲潄吀瑩敬')
        if corruption_chars / len(text_sample) > 0.25:  # 降低阈值，减少误判
            logger.debug(f"检测到高频异常字符: {corruption_chars}/{len(text_sample)}")
            return True
        
        # 3. 检测重复模式（优化算法）
        import re
        # 更精确的重复模式检测
        repeated_patterns = re.findall(r'(.{2,4})\1{6,}', text_sample)  # 提高重复阈值
        if len(repeated_patterns) > 12:  # 提高重复模式阈值
            logger.debug(f"检测到大量重复模式: {len(repeated_patterns)}")
            return True
        
        # 4. 智能可读性分析
        readability_score = self._calculate_readability_score(text_sample)
        if readability_score < 0.08:  # 进一步降低可读性阈值
            logger.debug(f"可读性评分过低: {readability_score}")
            return True
        
        # 5. 技术内容存在性检查
        has_technical_content = self._has_technical_content(text_sample)
        if not has_technical_content and readability_score < 0.15:
            logger.debug("无技术内容且可读性差")
            return True
        
        logger.debug(f"✅ 增强文本质量检测通过 - 编码: {encoding_score:.3f}, 可读性: {readability_score:.3f}")
        return False
    
    def _evaluate_text_quality(self, text: str) -> float:
        """评估提取文本的质量（0-1分数）"""
        if not text or len(text.strip()) < self.min_text_length:
            return 0.0
        
        text_clean = text.strip()
        
        # 🔍 乱码检测 - 如果检测到严重乱码，直接返回0
        if self._is_corrupted_text(text_clean):
            logger.warning(f"检测到乱码文本，质量评分为0")
            return 0.0
        
        score = 0.0
        
        # 1. 基础长度分数 (0.3权重)
        if len(text_clean) > 100:
            length_score = min(len(text_clean) / 1000, 1.0) * 0.3
            score += length_score
        
        # 2. 中文字符比例 (0.2权重)
        chinese_chars = len([c for c in text_clean if '\u4e00' <= c <= '\u9fff'])
        if len(text_clean) > 0:
            chinese_ratio = chinese_chars / len(text_clean)
            # 电气设备文档通常包含中文
            if chinese_ratio > 0.1:
                score += chinese_ratio * 0.2
        
        # 3. 技术关键词密度 (0.25权重)
        technical_keywords = [
            '产品', '型号', '规格', '参数', '技术', '电压', '电流', '功率', 
            '频率', '测试', '保护', '继电', '装置', '设备', '性能', '精度',
            'product', 'model', 'specification', 'parameter', 'voltage', 
            'current', 'power', 'frequency', 'test', 'protection'
        ]
        
        keyword_count = 0
        text_lower = text_clean.lower()
        for keyword in technical_keywords:
            keyword_count += text_lower.count(keyword.lower())
        
        keyword_density = min(keyword_count / max(len(text_clean.split()), 1), 0.1) * 10
        score += keyword_density * 0.25
        
        # 4. 数字和单位存在性 (0.15权重)
        import re
        number_pattern = r'\d+\.?\d*\s*[A-Za-z/Ω%°℃V]+'
        number_matches = len(re.findall(number_pattern, text_clean))
        if number_matches > 0:
            number_score = min(number_matches / 20, 1.0) * 0.15
            score += number_score
        
        # 5. 可读性检查 (0.1权重)
        # 检查是否包含过多乱码字符
        printable_chars = len([c for c in text_clean if c.isprintable() or c in '\n\t'])
        if len(text_clean) > 0:
            readability = printable_chars / len(text_clean)
            if readability > 0.8:
                score += 0.1
            elif readability > 0.6:
                score += 0.05
        
        # 确保分数在0-1范围内
        final_score = min(max(score, 0.0), 1.0)
        
        return final_score
    
    def _evaluate_ocr_quality(self, text: str) -> float:
        """评估OCR提取文本的质量（0-1分数）"""
        if not text or len(text.strip()) < 5:
            return 0.0
        
        text_clean = text.strip()
        score = 0.0
        
        # 1. 基础长度分数 (0.2权重)
        if len(text_clean) > 20:
            length_score = min(len(text_clean) / 200, 1.0) * 0.2
            score += length_score
        
        # 2. 可读字符比例 (0.3权重)
        readable_chars = len([c for c in text_clean if c.isalnum() or c.isspace() or c in '.,;:!?-()[]{}'])
        if len(text_clean) > 0:
            readable_ratio = readable_chars / len(text_clean)
            if readable_ratio > 0.7:
                score += 0.3
            elif readable_ratio > 0.5:
                score += 0.2
            elif readable_ratio > 0.3:
                score += 0.1
        
        # 3. 中文或英文字符存在性 (0.25权重)
        chinese_chars = len([c for c in text_clean if '\u4e00' <= c <= '\u9fff'])
        english_words = len([word for word in text_clean.split() if word.isalpha() and len(word) >= 3])
        
        if chinese_chars > 0 or english_words > 0:
            content_score = min((chinese_chars + english_words * 2) / 50, 1.0) * 0.25
            score += content_score
        
        # 4. 数字和技术特征 (0.15权重)
        import re
        has_numbers = bool(re.search(r'\d', text_clean))
        has_units = bool(re.search(r'\d+\s*[A-Za-z%°℃Ω]', text_clean))
        
        if has_numbers:
            score += 0.1
        if has_units:
            score += 0.05
        
        # 5. 避免明显的OCR错误 (0.1权重)
        # 检查是否有过多的单个字符或明显错误
        single_chars = len([word for word in text_clean.split() if len(word) == 1])
        total_words = len(text_clean.split())
        
        if total_words > 0:
            single_char_ratio = single_chars / total_words
            if single_char_ratio < 0.3:  # 单字符比例低于30%
                score += 0.1
        
        return min(max(score, 0.0), 1.0)
    
    def _contains_garbage_chinese(self, text: str) -> bool:
        """检测中文文本是否包含乱码汉字"""
        # 常见的.doc解析错误产生的乱码汉字
        garbage_chinese_chars = set([
            '潗', '摲', '楍', '牣', '獯', '景', '煅', '慵', '楴', '湯', '畱', '瑡', '潩',
            '卍', '潗', '摲', '潄', '吀', '瑩', '敬', '牁', '慩', '袈', '霡', '蠈', '袢',
            '屜', '屝', '屬', '屭', '屨', '屪', '屢', '屣', '層', '履', '屦', '屧', '屨', '屩', '屲'
        ])
        
        garbage_count = sum(1 for c in text if c in garbage_chinese_chars)
        garbage_ratio = garbage_count / len(text) if text else 0
        
        # 如果垃圾汉字超过20%，认为是乱码
        return garbage_ratio > 0.2
    
    def _is_meaningful_english(self, text: str) -> bool:
        """检测英文文本是否有意义（非乱码）"""
        import re
        
        # 检查是否包含常见的英文单词模式
        words = text.lower().split()
        
        # 技术文档常见的英文单词
        common_technical_words = set([
            'the', 'and', 'of', 'to', 'a', 'in', 'for', 'is', 'on', 'with', 'as', 'by',
            'test', 'voltage', 'current', 'power', 'frequency', 'protection', 'relay',
            'device', 'equipment', 'specification', 'parameter', 'technical', 'model',
            'product', 'range', 'accuracy', 'measurement', 'control', 'system',
            'output', 'input', 'signal', 'data', 'interface', 'communication'
        ])
        
        meaningful_words = 0
        for word in words:
            clean_word = re.sub(r'[^a-z]', '', word)  # 移除标点
            if (len(clean_word) >= 3 and 
                (clean_word in common_technical_words or 
                 re.match(r'^[a-z]+$', clean_word))):  # 纯英文字母组成
                meaningful_words += 1
        
        # 如果有意义的单词比例超过60%，认为是有意义的英文
        return meaningful_words / max(len(words), 1) > 0.6
    
    def _is_corrupted_tech_data(self, text: str) -> bool:
        """检测技术数据是否损坏"""
        import re
        
        # 检查是否包含明显的乱码模式
        corruption_patterns = [
            r'[潗摲楍牣獯景煅慵楴湯畱瑡潩卍潗摲潄吀瑩敬牁慩袈霡蠈袢]',  # 乱码汉字
            r'[▉▊▋▌▍▎▏█▄▀■□▲△▼▽◆◇○●◎☆★]',  # 图形字符
            r'[\ue000-\uf8ff]',  # 私用区字符
            r'[�?]{2,}',         # 替代字符
        ]
        
        for pattern in corruption_patterns:
            if re.search(pattern, text):
                return True
        
        return False
    
    def _is_obvious_garbage_line(self, line: str) -> bool:
        """检测是否为明显的垃圾行"""
        import re
        
        if not line or len(line.strip()) < 2:
            return True
        
        line = line.strip()
        
        # 检测明显的垃圾模式
        garbage_patterns = [
            r'^[潗摲楍牣獯景煅慵楴湯畱瑡潩卍潗摲潄吀瑩敬牁慩袈霡蠈袢]+$',  # 全乱码汉字
            r'^[▉▊▋▌▍▎▏█▄▀■□▲△▼▽◆◇○●◎☆★]+$',  # 全图形字符
            r'^[\u0080-\u00ff]{2,}$',         # 全高位ASCII
            r'^[�?]{2,}$',                    # 全替代字符
            r'^(.)\1{5,}$',                   # 同一字符重复6次以上
            r'^[\s\-\+\=\|]{3,}$',           # 表格分隔符
        ]
        
        for pattern in garbage_patterns:
            if re.search(pattern, line):
                return True
        
        # 检测可读性：如果可读字符少于50%，认为是垃圾
        readable_chars = sum(1 for c in line if (c.isalnum() or c.isspace() or 
                                                c in '.,;:!?-()[]{}""''、。，；：！？'))
        readable_ratio = readable_chars / len(line)
        
        return readable_ratio < 0.5
    
    def _clean_ocr_text(self, text: str) -> str:
        """清理OCR提取的文本"""
        import re
        
        # 1. 移除明显的OCR错误字符
        # 移除孤立的奇怪字符
        cleaned = re.sub(r'[^\w\s\u4e00-\u9fff.,;:!?\-()[\]{}\'"°℃%Ω/\\]', ' ', text)
        
        # 2. 规范空白字符
        cleaned = re.sub(r'\s+', ' ', cleaned)
        cleaned = re.sub(r'\n\s*\n', '\n\n', cleaned)
        
        # 3. 移除过短的行（可能是OCR噪音）
        lines = cleaned.split('\n')
        filtered_lines = []
        
        for line in lines:
            line = line.strip()
            # 保留较长的行或包含数字/技术信息的行
            if len(line) >= 3 or re.search(r'\d|[A-Za-z]{3,}|[\u4e00-\u9fff]{2,}', line):
                filtered_lines.append(line)
        
        # 4. 重新组合
        result = '\n'.join(filtered_lines)
        
        # 5. 最终清理
        result = result.strip()
        
        return result